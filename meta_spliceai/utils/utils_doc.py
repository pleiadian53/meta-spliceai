# pip install python-pptx

import os, sys

import pandas as pd 
import numpy as np

from pathlib import Path
from rich.console import Console  # pip install rich 
from rich.panel import Panel

# Create an instance of Console
console = Console()


def print_with_indent_v0(message, indent_level=0, indent_size=4):
    indent = ' ' * (indent_level * indent_size)
    print(f"{indent}{message}")


def print_with_indent(message, indent_level=0, indent_size=4):
    indent = ' ' * (indent_level * indent_size)
    console.print(f"{indent}{message}")


def print_section_separator():
    console.print(Panel("-" * 85, style="bold"))


def print_emphasized(text, style='bold', edge_effect=True, symbol='='):
    styles = {
        'bold': '\033[1m',
        'underline': '\033[4m',
        'red': '\033[91m',
        'green': '\033[92m',
        'yellow': '\033[93m',
        'blue': '\033[94m',
        'magenta': '\033[95m',
        'cyan': '\033[96m'
    }
    end_style = '\033[0m'  # Reset to default style

    if edge_effect:
        edge_line = symbol * len(text)
        print(edge_line)
        print(f"{styles.get(style, '')}{text}{end_style}")
        print(edge_line)
    else:
        print(f"{styles.get(style, '')}{text}{end_style}")

def calculate_table_dimension(df): 
    from pptx.util import Inches, Pt

    rows, cols = df.shape

    # Constants for table dimensions
    MAX_TABLE_WIDTH = Inches(8)  # max width that looks good on a slide
    MAX_TABLE_HEIGHT = Inches(5.5)  # max height considering space for a title and margins

    # Determine column width based on the number of columns, but ensure it doesn't exceed MAX_TABLE_WIDTH
    col_width = min(MAX_TABLE_WIDTH / cols, Inches(2))

    # Determine row height based on the number of rows (including header), but ensure it doesn't exceed MAX_TABLE_HEIGHT
    row_height = min(MAX_TABLE_HEIGHT / (rows + 1), Inches(0.5))

    # Calculate total table width and height
    table_width = col_width * cols
    table_height = row_height * (rows + 1)

    return table_width, table_height

def dataframe_to_ppt(df, **kargs): 
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Length

    verbose = kargs.get("verbose", 1)

    # Initialize a PowerPoint presentation
    prs = Presentation()

    # Add a slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add title to the slide
    title_text = kargs.get("title", "Classification Report")
    title = slide.shapes.title
    title.text = title_text # "Classification Report"

    # Define table dimensions
    rows, cols = df.shape
    left = Inches(1)
    top = Inches(2)  # Adjusted to accommodate title

    # width, height = calculate_table_dimension(df)
    width = Inches(8)
    height = Inches(4)

    # Add a table to the slide
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

    # Set column names as header row
    for col, col_name in enumerate(df.columns):
        cell = table.cell(0, col)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(91, 155, 213)
        cell.text_frame.paragraphs[0].font.size = Pt(14)  # Adjusted font size

    # Populate the table with DataFrame values
    for row in range(rows):
        for col in range(cols):

            # Check if the value is numeric and format it
            value = df.iloc[row, col]
            if isinstance(value, (int, float)):
                cell_text = "{:.3f}".format(value)
            else:
                cell_text = str(value)

            cell = table.cell(row + 1, col)
            cell.text = cell_text # str(df.iloc[row, col])
            
            cell.text_frame.paragraphs[0].font.size = Pt(12)  # Adjusted font size
            
            # Set alternating row colors for better readability
            if row % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(230, 230, 230)

    # Save the PowerPoint presentation
    output_dir_default = os.path.join(os.getcwd(), "data")
    output_dir = kargs.get("output_dir", output_dir_default)
    # Path(output_dir).mkdir(parents=True, exist_ok=True)

    ext = kargs.get("ext", "pptx")
    output_file = kargs.get("output_file", f"classification_report.{ext}")

    output_path = os.path.join(output_dir, output_file)
    if verbose: 
        print(f"[doc] Saving classification report slides to:\n{output_path}\n")
    # savefig(plt, output_path, ext=ext, dpi=100, message='', verbose=True)

    prs.save(output_path)

    return


def demo_dataframe_ppt(): 
    import pandas as pd
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor  # Updated import for RGBColor

    # Create a sample DataFrame
    data = {
        'Name': ['Alice', 'Bob', 'Charlie'],
        'Age': [25, 30, 28],
        'Country': ['USA', 'Canada', 'UK']
    }

    df = pd.DataFrame(data)

    # Create a PowerPoint presentation
    prs = Presentation()

    # Add a slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Define table dimensions
    # - calculate the dimensions of the table based on the number of rows and columns in the DataFrame.
    rows, cols = df.shape
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(4)

    # Add a table to the slide
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

    # Set column names as header row
    for col, col_name in enumerate(df.columns):
        cell = table.cell(0, col)
        cell.text = col_name
        cell.fill.solid()
        
        # Set header cell color using RGBColor
        cell.fill.fore_color.rgb = RGBColor(91, 155, 213)

    # Populate the table with DataFrame values
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row + 1, col)
            cell.text = str(df.iloc[row, col])

    # Save the PowerPoint presentation
    prs.save('dataframe_table.pptx')

    return


def test(): 

    # demo_dataframe_ppt() # ... ok
    
    data =  {
        'Name': ['Alice', 'Bob', 'Charlie'],
        'Age': [25, 30, 28],
        'Country': ['USA', 'Canada', 'UK']
    }
    df = pd.DataFrame(data)
    dataframe_to_ppt(df, title_text="Test Dataframe")

    return

if __name__ == "__main__": 
    test()
