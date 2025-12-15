import os
import numpy as np
from tabulate import tabulate
from typing import Union, List, Set

import pandas as pd
import polars as pl

from pybedtools import BedTool
from Bio import SeqIO

from .utils_bio import (
    normalize_strand
)

from .utils_doc import (
    print_emphasized, 
    print_with_indent, 
    print_section_separator, 
    display_dataframe_in_chunks
)

from .utils_df import (
    analyze_dataframe_properties
)

import pyBigWig
from tqdm import tqdm

# Refactor to utils_plot
import matplotlib.pyplot as plt 

from sklearn.metrics import RocCurveDisplay, PrecisionRecallDisplay
from sklearn.model_selection import StratifiedKFold
# from sklearn.model_selection import cross_val_predict, train_test_split
from sklearn.metrics import (
    precision_score, recall_score, f1_score, roc_auc_score,
    roc_curve, precision_recall_curve, auc
)

from meta_spliceai.splice_engine.model_evaluator import ModelEvaluationFileHandler
from typing import Optional


class PerformanceAnalyzer:
    eval_dir = '/path/to/meta-spliceai/data/ensembl/spliceai_eval'
    analysis_dir = '/path/to/meta-spliceai/data/ensembl/spliceai_analysis'

    # Todo: Add higher level modules or classes to initialize the FeatureAnalyzer
    #       For instance, both FeatureAanalyzer, ErrorAnalzyer can be initialized through a common class with 
    #       shared parameters (and methods and attributes)

    def __init__(self, output_dir=None):
        if output_dir is None: 
            output_dir = PerformanceAnalyzer.analysis_dir
        self.output_dir = output_dir

    def analyze_hard_genes(self, performance_df=None, *, metric='f1_score', error_type='FP', n_genes=100, save=False, **kargs): 
        
        strategy = kargs.get('strategy', 'worst')
        threshold = kargs.get('threshold', 0.5)
        
        # Call the global function
        df_hard_genes = analyze_hard_genes(
            performance_df, 
            metric=metric, 
            error_type=error_type, 
            n_genes=n_genes, 
            strategy=strategy, threshold=threshold)

        # Save the output if required
        if save:
            suffix = kargs.get('suffix') or error_type.lower() 
            output_path = os.path.join(self.output_dir, f'hard_genes_{suffix}.tsv')
            self.save_dataframe(df_hard_genes, output_path)
            print_emphasized(f"[i/o] Saved hard genes to {output_path}")

        return df_hard_genes

    def load_hard_genes(self, file_path=None, error_type='FP', raise_error=False):
        """
        Load a dataset of hard genes from a TSV file. Hard genes are those for which 
        SpliceAI has difficulty making accurate predictions.

        Parameters
        ----------
        file_path : str, optional
            Path to the TSV file containing hard genes data. If None, constructs path using
            self.output_dir and error_type (default: None)
        error_type : str, optional
            Type of prediction error to load ('FP' for False Positives or 'FN' for False Negatives).
            Used to construct the default file path if file_path is None (default: 'FP')
        raise_error : bool, optional
            If True, raises FileNotFoundError when file is not found.
            If False, prints warning message and returns None (default: False)

        Returns
        -------
        polars.DataFrame or pandas.DataFrame or None
            DataFrame containing hard genes data if file exists, None if file doesn't exist
            and raise_error is False

        Raises
        ------
        FileNotFoundError
            If the file doesn't exist and raise_error is True

        Examples
        --------
        >>> analyzer = PerformanceAnalyzer()
        >>> # Load using default path for False Positives
        >>> df = analyzer.load_hard_genes()
        >>> # Load from specific path
        >>> df = analyzer.load_hard_genes('/path/to/hard_genes.tsv')
        >>> # Load False Negatives, raising error if file not found
        >>> df = analyzer.load_hard_genes(error_type='FN', raise_error=True)
        """

        if file_path is None: 
            suffix = error_type.lower()
            file_path = os.path.join(self.output_dir, f'hard_genes_{suffix}.tsv')

        if not os.path.exists(file_path): 
            msg = f"Hard genes file not found at {file_path}"
            if raise_error: 
                raise FileNotFoundError(msg)
            else: 
                print(msg)
                return None

        return self.load_dataframe(file_path)

    def retrieve_hard_genes(self, performance_df=None, metric='f1_score', error_type='FP', n_genes=100, save=False, **kargs):
        """
        Retrieve or analyze genes that are difficult for SpliceAI to predict accurately.
        If previously analyzed genes exist and overwrite=False, loads from file.
        Otherwise performs new analysis.

        Parameters
        ----------
        performance_df : pd.DataFrame or pl.DataFrame, optional
            DataFrame containing performance metrics for each gene. If None, will be loaded
            from default location (default: None)
        metric : str, optional
            Performance metric to use for identifying hard genes (e.g. 'f1_score', 'precision')
            (default: 'f1_score')
        error_type : str, optional
            Type of prediction error to analyze ('FP' for False Positives or 'FN' for False Negatives)
            (default: 'FP')
        n_genes : int, optional
            Number of hard genes to retrieve (default: 100)
        save : bool, optional
            Whether to save the analyzed hard genes to file (default: False)

        Other Parameters
        ---------------
        verbose : int, optional
            Level of verbosity in output (default: 0)
        overwrite : bool, optional
            If True, forces new analysis even if saved data exists (default: False)
        strategy : str, optional
            Strategy for selecting hard genes ('worst' or 'average') (default: 'worst')
        threshold : float, optional
            Performance threshold below which genes are considered "hard" (default: 0.5)

        Returns
        -------
        polars.DataFrame or pandas.DataFrame
            DataFrame containing the identified hard genes and their performance metrics

        Notes
        -----
        "Hard genes" are those for which SpliceAI has difficulty making accurate predictions,
        as measured by the specified performance metric being below the threshold and/or
        having a high number of prediction errors of the specified type.
        """
        verbose = kargs.get('verbose', 0)
        overwrite = kargs.get('overwrite', False)

        df = None
        if not overwrite:
            df = self.load_hard_genes(error_type=error_type)

        if df is None: 
            strategy = kargs.get('strategy', 'worst')
            threshold = kargs.get('threshold', 0.5)
            
            # "Hard genes" are those genes for which SpliceAI didn't do well
            print_emphasized(
                f"[workflow] Analyzing 'hard' genes with strategy={strategy}, threshold={threshold} ...")
            
            df = self.analyze_hard_genes(
                performance_df, metric=metric, error_type=error_type, 
                strategy=strategy, 
                threshold=threshold, 
                n_genes=n_genes,
                save=save, 
                verbose=verbose)
        
        return df


    def load_dataframe(self, file_path, separator=None):
        # Ensure the file_path is within self.output_dir if it doesn't have a preceding directory
        if not os.path.dirname(file_path):
            file_path = os.path.join(self.output_dir, file_path)

        # Determine the separator based on the file extension if not provided
        if separator is None:
            if file_path.endswith('.csv'):
                separator = ','
            elif file_path.endswith('.tsv'):
                separator = '\t'
            else:
                raise ValueError("Unsupported file format")

        try:
            return pl.read_csv(file_path, separator=separator)
        except Exception:
            return pd.read_csv(file_path, sep=separator)

    def save_dataframe(self, df, file_path, separator=None):
        # Ensure the file_path is within self.output_dir if it doesn't have a preceding directory
        if not os.path.dirname(file_path):
            file_path = os.path.join(self.output_dir, file_path)

        # Determine the separator based on the file extension if not provided
        if separator is None:
            if file_path.endswith('.csv'):
                separator = ','
            elif file_path.endswith('.tsv'):
                separator = '\t'
            else:
                raise ValueError("Unsupported file format")

        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        if isinstance(df, pl.DataFrame):
            df.write_csv(file_path, separator=separator)
        elif isinstance(df, pd.DataFrame):
            df.to_csv(file_path, sep=separator, index=False)
        else:
            raise ValueError("Unsupported DataFrame type")

        return file_path


# Function to compute summary statistics
def compute_summary_stats(data):
    summary = {
        'Total Genes': len(data['gene_id'].unique()),
        'Average Precision': data['precision'].mean(),
        'Average Recall': data['recall'].mean(),
        'Average Specificity': data['specificity'].mean(),
        'Average F1 Score': data['f1_score'].mean(),
        'Average False Positive Rate': data['fpr'].mean(),
        'Average False Negative Rate': data['fnr'].mean(),
        'Genes with Perfect Precision and Recall': len(data[(data['precision'] == 1.0) & (data['recall'] == 1.0)]),
        'Genes with Zero Precision or Recall': len(data[(data['precision'] == 0.0) | (data['recall'] == 0.0)]),
    }
    return summary


def analyze_performance_profile(
    performance_df=None, 
    metric_to_plot='f1_score', 
    save_plot=False, 
    plot_format='pdf', 
    plot_file='metric_distribution.pdf', 
    save_stats=False, 
    stats_file='summary_stats.csv', 
    verbose=True, 
    **kargs
):
    """
    Analyze SpliceAI performance and provide insights via statistics and visualizations.

    Parameters:
    - performance_df (pd.DataFrame): DataFrame containing SpliceAI performance data.
    - metric_to_plot (str): The performance metric to visualize (e.g., 'f1_score', 'precision').
    - save_plot (bool): Whether to save the plot (default is False).
    - plot_format (str): Format to save the plot (default is 'pdf').
    - plot_file (str): File name to save the plot (default is 'metric_distribution.pdf').
    - save_stats (bool): Whether to save the summary statistics (default is False).
    - stats_file (str): File name to save the summary statistics (default is 'summary_stats.csv').
    - verbose (bool): Whether to print progress updates (default is True).

    Returns:
    - None
    """
    import pandas as pd
    import matplotlib.pyplot as plt
    # from tabulate import tabulate
    # from meta_spliceai.splice_engine.model_evaluator import ModelEvaluationFileHandler

    if performance_df is None: 
        if verbose: 
            print("Loading performance data...")
        
        # Replace this with your data loading logic
        eval_dir = kargs.get('eval_dir', PerformanceAnalyzer.eval_dir)
        separator = kargs.get('separator', kargs.get('sep', '\t'))

        # Initialize ModelEvaluationFileHandler
        mefd = ModelEvaluationFileHandler(eval_dir, separator=separator)

        # Load the performance dataframe
        performance_df = mefd.load_performance_df(aggregated=True)

    # Ensure the DataFrame is in Pandas format for compatibility
    if isinstance(performance_df, pl.DataFrame):
        performance_df = performance_df.to_pandas()

    if 'n_splice_sites' not in performance_df.columns:
        raise ValueError("The input DataFrame must contain 'n_splice_sites' column.")

    total_genes = len(performance_df['gene_id'].unique())

    # Filter out rows with None or null values
    performance_df = performance_df.dropna()

    # Filter out genes with no valid splice sites
    valid_data = performance_df[performance_df['n_splice_sites'] > 0]
    if valid_data.empty:
        raise ValueError("No valid genes found after filtering. Please check your data.")

    n_genes_without_splice_sites = performance_df[performance_df['n_splice_sites'] == 0]
    if verbose: 
        print(f"Number of genes without splice sites: {len(n_genes_without_splice_sites)}")

    # Compute summary statistics
    valid_summary_stats = {
        'Total Genes': len(valid_data['gene_id'].unique()),
        'Average Precision': valid_data['precision'].mean(),
        'Average Recall': valid_data['recall'].mean(),
        'Average Specificity': valid_data['specificity'].mean(),
        'Average F1 Score': valid_data['f1_score'].mean(),
        'Average False Positive Rate': valid_data['fpr'].mean(),
        'Average False Negative Rate': valid_data['fnr'].mean(),
        'Genes with Perfect Precision and Recall': len(valid_data[(valid_data['precision'] == 1.0) & (valid_data['recall'] == 1.0)]),
        'Genes with Zero Precision or Recall': len(valid_data[(valid_data['precision'] == 0.0) | (valid_data['recall'] == 0.0)]),
        'Perfect Precision and Recall Rate': len(valid_data[(valid_data['precision'] == 1.0) & (valid_data['recall'] == 1.0)]) / total_genes,
        'Zero Precision or Recall Rate': len(valid_data[(valid_data['precision'] == 0.0) | (valid_data['recall'] == 0.0)]) / total_genes,
    }

    # Print summary statistics in a readable format
    if verbose:
        print("\nFiltered SpliceAI Performance Summary Statistics:")
        print(tabulate(pd.DataFrame.from_dict(valid_summary_stats, orient='index', columns=['Value']), headers="keys", tablefmt="pretty"))

    # Save summary statistics if needed
    if save_stats:
        output_dir = kargs.get('output_dir', PerformanceAnalyzer.analysis_dir)
        stats_path = os.path.join(output_dir, stats_file)
        pd.DataFrame.from_dict(valid_summary_stats, orient='index', columns=['Value']).to_csv(stats_path)

        if verbose: 
            print(f"Summary statistics saved to {stats_path}")

        pptx_file_path = os.path.join(output_dir, 'summary_stats.pptx')
        save_summary_stats_to_pptx(
            valid_summary_stats, 
            pptx_file=pptx_file_path, 
            slide_title='Filtered SpliceAI Performance Summary Statistics')

        if verbose: 
            print(f"Summary statistics (slides) saved to {pptx_file_path}")

    # Ensure the metric_to_plot exists in the DataFrame
    if metric_to_plot not in valid_data.columns:
        raise ValueError(f"The specified metric '{metric_to_plot}' is not in the DataFrame. Available metrics: {list(valid_data.columns)}")

    # Plot the selected metric distribution
    fig = plt.figure(figsize=(8, 6))
    plt.hist(valid_data[metric_to_plot].dropna(), bins=20, edgecolor='k', alpha=0.7)
    plt.title(f"Distribution of {metric_to_plot.capitalize()} (Filtered for Valid Splice Sites)")
    plt.xlabel(metric_to_plot.capitalize())
    plt.ylabel("Frequency")

    if save_plot:
        output_dir = kargs.get('output_dir', PerformanceAnalyzer.analysis_dir)
        output_path = os.path.join(output_dir, plot_file)
        plt.savefig(output_path, format=plot_format)
        if verbose: 
            print(f"Plot saved as {output_path}")
    else:
        plt.show()

    plt.close(fig)


def sort_performance_by_metric_and_error_type(
        df, metric='f1_score', error_type='FP', df_type='polars', 
        min_splice_sites=0, splice_type=None, aggregate='average', primary_threshold=0.5):
    """
    Sort the rows of the DataFrame by a specified performance metric in ascending order
    and then by error type (either FP or FN).

    Parameters:
    - df (pl.DataFrame or pd.DataFrame): The input DataFrame.
    - metric (str): The performance metric to sort by (default is 'f1_score').
    - error_type (str): The error type to sort by (either 'FP' or 'FN', default is 'FP').
    - df_type (str): The type of DataFrame ('polars' or 'pandas', default is 'polars').
    - min_splice_sites (int): Minimum number of splice sites to include (default is 0).
    - splice_type (str): The splice type to filter by (either 'donor' or 'acceptor', default is None).
    - aggregate (str): The aggregation method to use ('average', 'max', 'min', default is 'average').
    - primary_threshold (float): The threshold for the primary sorting criterion, i.e. performance metric like F1 score 
                                (default is 0.5).

    Returns:
    - pl.DataFrame or pd.DataFrame: The sorted DataFrame.
    """
    if metric not in df.columns:
        raise ValueError(f"Metric '{metric}' not found in DataFrame columns.")
    if error_type not in ['FP', 'FN']:
        raise ValueError(f"Error type '{error_type}' must be either 'FP' or 'FN'.")

    # Ensure 'n_splice_sites' and 'splice_type' columns are present
    if 'n_splice_sites' not in df.columns or 'splice_type' not in df.columns:
        raise ValueError("The input DataFrame must contain both 'n_splice_sites' and 'splice_type' columns.")

    if df_type == 'polars':
        # Filter out rows where 'n_splice_sites' is too small or 0
        df = df.filter(pl.col('n_splice_sites') > min_splice_sites)

        # Filter by 'splice_type' if it is specified
        if splice_type is not None:
            df = df.filter(pl.col('splice_type') == splice_type)

        # Group by 'gene_id' and aggregate the performance metric
        if aggregate == 'average':
            grouped_df = df.group_by('gene_id').agg([
                pl.col(metric).mean().alias(f'{metric}_agg'),
                pl.col(error_type).sum().alias(error_type)
            ])
        elif aggregate == 'max':
            grouped_df = df.group_by('gene_id').agg([
                pl.col(metric).max().alias(f'{metric}_agg'),
                pl.col(error_type).sum().alias(error_type)
            ])
        elif aggregate == 'min':
            grouped_df = df.group_by('gene_id').agg([
                pl.col(metric).min().alias(f'{metric}_agg'),
                pl.col(error_type).sum().alias(error_type)
            ])
        else:
            raise ValueError(f"Unsupported aggregation method '{aggregate}'. Use 'average', 'max', or 'min'.")

        # Filter by the primary sorting criterion
        primary_filtered_df = grouped_df.filter(pl.col(f'{metric}_agg') <= primary_threshold)

        # Sort by the primary sorting criterion in ascending order
        sorted_df = primary_filtered_df.sort(f'{metric}_agg', descending=False)

        # Further sort by the secondary sorting criterion in descending order
        sorted_df = sorted_df.sort(error_type, descending=True)

    elif df_type == 'pandas':
        # Filter out rows where 'n_splice_sites' is too small or 0
        df = df[df['n_splice_sites'] > min_splice_sites]

        # Filter by 'splice_type' if it is specified
        if splice_type is not None:
            df = df[df['splice_type'] == splice_type]

        # Group by 'gene_id' and aggregate the performance metric
        if aggregate == 'average':
            grouped_df = df.groupby('gene_id').agg({
                metric: 'mean',
                error_type: 'sum'
            }).reset_index().rename(columns={metric: f'{metric}_agg'})
        elif aggregate == 'max':
            grouped_df = df.groupby('gene_id').agg({
                metric: 'max',
                error_type: 'sum'
            }).reset_index().rename(columns={metric: f'{metric}_agg'})
        elif aggregate == 'min':
            grouped_df = df.groupby('gene_id').agg({
                metric: 'min',
                error_type: 'sum'
            }).reset_index().rename(columns={metric: f'{metric}_agg'})
        else:
            raise ValueError(f"Unsupported aggregation method '{aggregate}'. Use 'average', 'max', or 'min'.")

        # Filter by the primary sorting criterion
        primary_filtered_df = grouped_df[grouped_df[f'{metric}_agg'] <= primary_threshold]

        # Sort by the primary sorting criterion in ascending order
        sorted_df = primary_filtered_df.sort_values(by=f'{metric}_agg', ascending=True)

        # Further sort by the secondary sorting criterion in descending order
        sorted_df = sorted_df.sort_values(by=error_type, ascending=False)

    else:
        raise ValueError(f"Unsupported DataFrame type '{df_type}'. Use 'polars' or 'pandas'.")

    return sorted_df


def select_worst_performing_genes(
    df, 
    metric='f1_score', 
    error_type='FP', 
    df_type='polars', 
    N=10, 
    strategy='worst',
    primary_threshold=0.5
):
    """
    Select the N unique genes for which the performance metric is worst.

    Parameters:
    - df (pl.DataFrame or pd.DataFrame): The input DataFrame.
    - metric (str): The performance metric to sort by (default is 'f1_score').
    - error_type (str): The error type to sort by (either 'FP' or 'FN', default is 'FP').
    - df_type (str): The type of DataFrame ('polars' or 'pandas', default is 'polars').
    - N (int): The number of unique genes to select (default is 10).
    - strategy (str): The strategy to use for selecting genes ('worst' or 'average', default is 'worst').
    - primary_threshold (float): The threshold for the primary sorting criterion (default is 0.5).

    Returns:
    - pl.DataFrame or pd.DataFrame: The DataFrame containing the N worst performing unique genes.
    """
    if df_type == 'polars':
        # Filter out rows where 'n_splice_sites' is too small or 0
        df = df.filter(pl.col('n_splice_sites') > 0)

        # Group by 'gene_id' and aggregate the performance metric and error type count
        if strategy == 'worst':
            print(f"[info] Using 'worst' strategy, at threshold={primary_threshold}")
            grouped_df = df.group_by('gene_id').agg([
                pl.col(metric).max().alias(f'{metric}_agg'),
                pl.col(error_type).sum().alias(f'{error_type}_count')
            ])
        elif strategy == 'average':
            print(f"[info] Using 'average' strategy, at threshold={primary_threshold}")
            grouped_df = df.group_by('gene_id').agg([
                pl.col(metric).mean().alias(f'{metric}_agg'),
                pl.col(error_type).sum().alias(f'{error_type}_count')
            ])
        else:
            raise ValueError(f"Unsupported strategy '{strategy}'. Use 'worst' or 'average'.")

        # Filter by the primary sorting criterion
        primary_filtered_df = grouped_df.filter(pl.col(f'{metric}_agg') <= primary_threshold)

        # Check if the number of genes is already smaller than or equal to N
        if primary_filtered_df.shape[0] <= N:
            return primary_filtered_df

        # Sort by the primary sorting criterion in ascending order
        sorted_df = primary_filtered_df.sort(f'{metric}_agg', descending=False)

        # Further sort by the secondary sorting criterion in descending order
        sorted_df = sorted_df.sort(f'{error_type}_count', descending=True)

        # Select the N unique genes with the worst performance
        worst_genes_df = sorted_df.head(N)

    elif df_type == 'pandas':
        # Filter out rows where 'n_splice_sites' is too small or 0
        df = df[df['n_splice_sites'] > 0]

        # Group by 'gene_id' and aggregate the performance metric and error type count
        if strategy == 'worst':
            print(f"[info] Using 'worst' strategy, at threshold={primary_threshold}")
            grouped_df = df.groupby('gene_id').agg({
                metric: 'max',
                error_type: 'sum'
            }).rename(columns={metric: f'{metric}_agg', error_type: f'{error_type}_count'})
        elif strategy == 'average':
            print(f"[info] Using 'average' strategy, at threshold={primary_threshold}")
            grouped_df = df.groupby('gene_id').agg({
                metric: 'mean',
                error_type: 'sum'
            }).rename(columns={metric: f'{metric}_agg', error_type: f'{error_type}_count'})
        else:
            raise ValueError(f"Unsupported strategy '{strategy}'. Use 'worst' or 'average'.")

        # Filter by the primary sorting criterion
        primary_filtered_df = grouped_df[grouped_df[f'{metric}_agg'] <= primary_threshold]

        # Check if the number of genes is already smaller than or equal to N
        if primary_filtered_df.shape[0] <= N:
            return primary_filtered_df

        # Sort by the primary sorting criterion in ascending order
        sorted_df = primary_filtered_df.sort_values(by=f'{metric}_agg', ascending=True)

        # Further sort by the secondary sorting criterion in descending order
        sorted_df = sorted_df.sort_values(by=f'{error_type}_count', ascending=False)

        # Select the N unique genes with the worst performance
        worst_genes_df = sorted_df.head(N)

    else:
        raise ValueError(f"Unsupported df_type '{df_type}'. Use 'polars' or 'pandas'.")

    return worst_genes_df


def save_summary_stats_to_pptx(summary_stats, pptx_file='summary_stats.pptx', slide_title='Summary Statistics'):
    """
    Save summary statistics to a PowerPoint slide.

    Parameters:
    - summary_stats (dict): Summary statistics to save.
    - pptx_file (str): File name to save the PowerPoint presentation (default is 'summary_stats.pptx').
    - slide_title (str): Title of the slide (default is 'Summary Statistics').

    Returns:
    - None
    """
    from pptx import Presentation
    from pptx.util import Inches

    # Create a presentation object
    prs = Presentation()

    # Add a slide with a title and content layout
    slide_layout = prs.slide_layouts[5]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title to the slide
    title = slide.shapes.title
    title.text = slide_title

    # Convert summary statistics to a DataFrame
    df = pd.DataFrame.from_dict(summary_stats, orient='index', columns=['Value'])

    # Add a table to the slide
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols + 1, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8 * (rows + 1))).table

    # Set column headers
    table.cell(0, 0).text = 'Metric'
    table.cell(0, 1).text = 'Value'

    # Populate the table with data
    for i, (index, row) in enumerate(df.iterrows()):
        table.cell(i + 1, 0).text = str(index)
        table.cell(i + 1, 1).text = str(row['Value'])

    # Save the presentation
    prs.save(pptx_file)
    print(f"Summary statistics saved to {pptx_file}")


def demo_sort_performance_by_metric_and_error_type(**kargs): 
    from .performance_analyzer import PerformanceAnalyzer

    print("Loading performance data...")
    
    # Replace this with your data loading logic
    eval_dir = kargs.get('eval_dir', PerformanceAnalyzer.eval_dir)
    separator = kargs.get('separator', kargs.get('sep', '\t'))

    # Initialize ModelEvaluationFileHandler
    mefd = ModelEvaluationFileHandler(eval_dir, separator=separator)

    # Load the performance dataframe
    performance_df = mefd.load_performance_df(aggregated=True)

    # Example usage
    sorted_full_performance_df = sort_performance_by_metric_and_error_type(performance_df, metric='f1_score', error_type='FP', df_type='polars')
    display_dataframe_in_chunks(sorted_full_performance_df, num_rows=50, title="Sorted Performance DataFrame")

    df_hard_genes = select_worst_performing_genes(performance_df, metric='f1_score', error_type='FP', df_type='polars', N=50)
    display_dataframe_in_chunks(df_hard_genes, num_rows=100, title="Worst Performing Genes")


def analyze_hard_genes(performance_df=None, *, metric='f1_score', error_type='FP', n_genes=100, **kargs): 
    from .extract_genomic_features import SpliceAnalyzer
    # from .performance_analyzer import PerformanceAnalyzer

    eval_dir = kargs.get('eval_dir', PerformanceAnalyzer.eval_dir)
    separator = kargs.get('separator', kargs.get('sep', '\t'))
    df_type = kargs.get('df_type', 'polars')
    verbose = kargs.get('verbose', 1)
    strategy = kargs.get('strategy', 'worst')
    threshold = kargs.get('threshold', 0.5)

    mefd = ModelEvaluationFileHandler(eval_dir, separator=separator)  # Initialize ModelEvaluationFileHandler

    if performance_df is None:
        print_emphasized("[i/o] Loading performance data...")
        performance_df = mefd.load_performance_df(aggregated=True)  # Load the performance dataframe

    df_hard_genes = \
        select_worst_performing_genes(
            performance_df,
            metric=metric, 
            error_type=error_type, 
            df_type='polars', N=n_genes, 
            strategy=strategy, 
            primary_threshold=threshold)
    display_dataframe_in_chunks(df_hard_genes, num_rows=100, title="Worst Performing Genes")

    # Overlapping genes
    if verbose > 0:
        print_emphasized("[i/o] Retrieving overlapping gene metadata...")
        sa = SpliceAnalyzer()
        overlapping_genes_metadata = sa.retrieve_overlapping_gene_metadata()

        # For each (hard) gene, look into its overlapping gene profile
        print_with_indent("Analyzing overlapping genes for hard genes ...", indent_level=1)
        df_hard_genes = df_hard_genes.to_pandas()
        for gene_id in df_hard_genes['gene_id'].unique(): 
            # gene_data = df_hard_genes[df_hard_genes['gene_id'] == gene_id].to_dict(orient='records')[0]
            # retrieve_overlapping_genes(gene_id, gene_data, overlapping_genes_metadata)

            n_overlapped_genes = 0
            if gene_id in overlapping_genes_metadata: 
                overlaps = overlapping_genes_metadata[gene_id]
                n_overlapped_genes = len(overlaps)
                print(f"Gene {gene_id} has {len(overlaps)} overlapping genes.")

                if n_overlapped_genes > 0: 
                    for i, overlap in enumerate(overlaps): 
                        print(f"  {i+1}. {overlap}")
            else: 
                if verbose > 1: 
                    print(f"Gene {gene_id} is NOT in the overlapping gene metadata.")

        print_with_indent("Analyzing overlapping genes for customized genes ...", indent_level=1)
        customized_genes = ['ENSG00000260597', ]
        for gene_id in customized_genes: 
            
            n_overlapped_genes = 0
            if gene_id in overlapping_genes_metadata: 
                overlaps = overlapping_genes_metadata[gene_id]
                n_overlapped_genes = len(overlaps)
                if verbose > 1: 
                    print(f"Gene {gene_id} has {len(overlaps)} overlapping genes.")

                    if n_overlapped_genes > 0: 
                        for i, overlap in enumerate(overlaps): 
                            print(f"  {i+1}. {overlap}")
            else: 
                if verbose > 1: 
                    print(f"Gene {gene_id} is NOT in the overlapping gene metadata.")

    if df_type == 'polars': 
        if not isinstance(df_hard_genes, pl.DataFrame):
            df_hard_genes = pl.from_pandas(df_hard_genes)
        return df_hard_genes

    return df_hard_genes

############################################
# from sklearn.metrics import RocCurveDisplay, PrecisionRecallDisplay


def plot_single_roc_curve(y_test, y_proba, subject="fp_vs_tp", output_path=None, save=True):
    plt.figure(figsize=(8, 6))
    disp = RocCurveDisplay.from_predictions(y_test, y_proba)
    # or disp = RocCurveDisplay.from_estimator(model, X_test, y_test)

    disp.ax_.set_title(f"{subject} - ROC Curve")
    disp.ax_.grid()

    if save and output_path:
        plt.savefig(output_path, bbox_inches="tight")
        print(f"[i/o] Saved ROC to: {output_path}")
        plt.close()
    else:
        plt.show()

def plot_single_pr_curve(y_test, y_proba, subject="fp_vs_tp", output_path=None, save=True):
    plt.figure(figsize=(8, 6))
    disp = PrecisionRecallDisplay.from_predictions(y_test, y_proba)

    disp.ax_.set_title(f"{subject} - PR Curve")
    disp.ax_.grid()

    if save and output_path:
        # output_file = os.path.join(output_dir, f"{subject}-prc-sklearn.pdf")
        plt.savefig(output_path, bbox_inches="tight")
        print(f"[i/o] Saved PR to: {output_path}")
        plt.close()
    else:
        plt.show()


def plot_cv_roc_curve(
    model,
    X, y,
    cv_splits=5,
    random_state=42,
    title="Cross-Validated ROC",
    output_path=None,
    show_std=True,
    plot_folds=False,   # Whether to plot individual fold ROC curves
    **fit_params
):
    """
    Plot a cross-validated ROC curve with mean +/- std shading.

    Important Subtlety on "Mean ROC / Mean AUC":
    -------------------------------------------
    There are two common ways to report a "Mean AUC" from cross-validation:

    (A) AUC of the *average TPR* curve:
        mean_tpr = average(interpolated TPRs across folds)
        mean_auc = auc(mean_fpr, mean_tpr)
        -> This represents the area under *one* "mean" curve.

    (B) Average of the per-fold AUCs:
        aucs = [ auc(fpr_i, tpr_i) for each fold i ]
        mean_auc = np.mean(aucs)
        std_auc  = np.std(aucs)
        -> This represents the *average* of the fold-specific AUC values.

    They are not always the same number! 
    Often people report (B) in the legend, but still plot (A) as the "mean ROC curve."
    """
    from sklearn.model_selection import StratifiedKFold
    from sklearn.metrics import roc_curve, auc
    from sklearn.base import clone

    # Set up cross-validation
    skf = StratifiedKFold(
        n_splits=cv_splits,
        shuffle=True,
        random_state=random_state
    )

    # A common FPR grid
    mean_fpr = np.linspace(0, 1, 100)
    tprs = []
    aucs = []

    fig = plt.figure(figsize=(8, 6))

    for fold_idx, (train_idx, test_idx) in enumerate(skf.split(X, y)):
        X_train, X_test = X.iloc[train_idx], X.iloc[test_idx]
        y_train, y_test = y.iloc[train_idx], y.iloc[test_idx]

        # Clone to avoid any carry-over state
        model_clone = clone(model)
        model_clone.fit(X_train, y_train, **fit_params)

        y_proba = model_clone.predict_proba(X_test)[:, 1]
        fpr, tpr, _ = roc_curve(y_test, y_proba)
        fold_auc = auc(fpr, tpr)
        aucs.append(fold_auc)

        if plot_folds:
            plt.plot(
                fpr, tpr, lw=1, alpha=0.3,
                label=f"Fold {fold_idx+1} (AUC={fold_auc:.2f})"
            )

        # Interpolate TPR at mean_fpr
        tpr_interp = np.interp(mean_fpr, fpr, tpr)
        tpr_interp[0] = 0.0
        tprs.append(tpr_interp)

    # Compute the mean TPR across folds
    mean_tpr = np.mean(tprs, axis=0)
    mean_tpr[-1] = 1.0  # Force the curve to end at TPR=1.0

    # --------------------------------------------------
    # Currently, the code does (A): AUC of the average TPR
    # --------------------------------------------------
    mean_auc_by_curve = auc(mean_fpr, mean_tpr)
    # This is the area under the single "mean" TPR curve
    # which is *not* necessarily the same as average of fold AUCs.

    # If you want to label with the *average of per-fold AUCs*, do:
    mean_auc_folds = np.mean(aucs)  # (B)
    std_auc_folds  = np.std(aucs)

    # For illustration, let's pick how we want to label the plot:
    # Suppose we want to show the average of fold AUCs (B) in the legend:
    # (That's more common for "cross-validated performance" reporting.)
    # We can store the "AUC of the mean curve" separately if we want.
    # --------------------------------------------------

    # Plot the "mean" ROC curve (visually from TPR interpolation)
    plt.plot(
        mean_fpr,
        mean_tpr,
        color='b',
        # NOTE: Using the average of fold AUCs in the legend:
        label=f"Mean ROC (AUC={mean_auc_folds:.2f} ± {std_auc_folds:.2f})",
        lw=2
    )

    # Optionally fill ±1 std around the mean TPR
    if show_std:
        std_tpr = np.std(tprs, axis=0)
        tpr_upper = np.minimum(mean_tpr + std_tpr, 1)
        tpr_lower = np.maximum(mean_tpr - std_tpr, 0)
        plt.fill_between(
            mean_fpr,
            tpr_lower,
            tpr_upper,
            color='grey',
            alpha=.2,
            label="± 1 std. dev."
        )

    plt.plot([0, 1], [0, 1], linestyle='--', color='r', label='Random Guess')

    plt.xlim([0, 1])
    plt.ylim([0, 1])
    plt.xlabel("False Positive Rate")
    plt.ylabel("True Positive Rate")
    plt.title(title)
    plt.legend(loc="lower right")
    plt.grid()

    # Save or show
    if output_path:
        plt.savefig(output_path, bbox_inches='tight')
        plt.close(fig)
        print(f"[i/o] Saved CV ROC curve to: {output_path}")
    else:
        plt.show()


def plot_cv_pr_curve(
    model,
    X, y,
    cv_splits=5,
    random_state=42,
    title="Cross-Validated PR Curve",
    output_path=None,
    show_std=True,
    plot_folds=False,
    **fit_params
):
    from sklearn.base import clone
    from sklearn.model_selection import StratifiedKFold
    from sklearn.metrics import precision_recall_curve, auc

    skf = StratifiedKFold(n_splits=cv_splits, shuffle=True, random_state=random_state)
    mean_recall = np.linspace(0, 1, 100)
    precisions = []
    aucs = []

    fig = plt.figure(figsize=(8, 6))

    for fold_idx, (train_idx, test_idx) in enumerate(skf.split(X, y)):
        X_train, X_test = X.iloc[train_idx], X.iloc[test_idx]
        y_train, y_test = y.iloc[train_idx], y.iloc[test_idx]

        # Clone the original model so that it is a fresh, untrained copy
        model_clone = clone(model)

        # Now fit the cloned model
        model_clone.fit(X_train, y_train, **fit_params)
        # model.fit(X_train, y_train, **fit_params)

        y_proba = model_clone.predict_proba(X_test)[:, 1]

        prec, rec, _ = precision_recall_curve(y_test, y_proba)
        pr_auc = auc(rec, prec)
        aucs.append(pr_auc)

        # Reverse if recall is descending
        if rec[0] > rec[-1]:
            prec, rec = prec[::-1], rec[::-1]

        # Interpolate at mean_recall
        prec_interp = np.interp(mean_recall, rec, prec)
        precisions.append(prec_interp)

        # Optionally plot each fold's curve
        if plot_folds:
            plt.plot(rec, prec, lw=1, alpha=0.3,
                     label=f"Fold {fold_idx+1} (AUC={pr_auc:.2f})")

    mean_prec = np.mean(precisions, axis=0)
    mean_auc = np.mean(aucs)
    std_auc = np.std(aucs)

    plt.plot(
        mean_recall,
        mean_prec,
        color='b',
        label=f"Mean PR (AUC={mean_auc:.2f} ± {std_auc:.2f})",
        lw=2
    )

    if show_std:
        std_prec = np.std(precisions, axis=0)
        prec_upper = np.minimum(mean_prec + std_prec, 1)
        prec_lower = np.maximum(mean_prec - std_prec, 0)
        plt.fill_between(
            mean_recall,
            prec_lower,
            prec_upper,
            color='grey',
            alpha=0.2,
            label="±1 std. dev."
        )

    plt.xlim([0, 1])
    plt.ylim([0, 1])
    plt.xlabel("Recall")
    plt.ylabel("Precision")
    plt.title(title)
    plt.grid()

    # Let matplotlib automatically use the labels defined above
    plt.legend(loc="best")

    if output_path:
        plt.savefig(output_path, bbox_inches='tight')
        plt.close(fig)
        print(f"[i/o] Saved CV PR curve to: {output_path}")
    else:
        plt.show()


###########################################
# from typing import Optional

def plot_gene_performance_2d(
    perf_df,
    x_col="recall",
    y_col="precision",
    plot_type="density",
    bins=50,
    figsize=(8,6),
    title=None,
    output_path=None,
    show=False,
    verbose=1
):
    """
    Visualize per-gene splicing performance in a 2D space (e.g., recall vs. precision or TPR vs. FPR).
    
    Parameters
    ----------
    perf_df : pd.DataFrame or pl.DataFrame
        Must contain the columns `x_col` and `y_col` plus presumably "gene_id" and others.
        Example columns: ["TP","FP","FN","TN","precision","recall","fpr","fnr","gene_id",...]
    x_col : str
        The name of the column for the x-axis (e.g., "recall", "precision", "fpr", "fnr", etc.).
    y_col : str
        The name of the column for the y-axis.
    plot_type : str
        - "density" => scatter plot with 2D density (kde or hist2d).
        - "heatmap" => binned heatmap (like a 2D histogram).
    bins : int
        Number of bins if using "heatmap" or hist2d approach.
    figsize : tuple
        Figure size in inches.
    title : str or None
        Custom plot title. If None, one is constructed from x_col, y_col, and plot_type.
    output_path : str or None
        If given, save the figure to this path (e.g. "my_plot.pdf"). 
        File extension (pdf/png) determines format.
    show : bool
        If True, calls plt.show() at the end.
    verbose : int
        Verbosity level. If > 0, prints some debug info.

    Returns
    -------
    fig, ax : Matplotlib figure and axes objects.
    """
    import polars as pl
    import pandas as pd
    import matplotlib.pyplot as plt
    import seaborn as sns
    import numpy as np

    # Convert Polars DataFrame to Pandas if necessary
    if isinstance(perf_df, pl.DataFrame):
        perf_df = perf_df.to_pandas()

    # Basic checks
    if x_col not in perf_df.columns:
        raise ValueError(f"Column '{x_col}' not found in perf_df.")
    if y_col not in perf_df.columns:
        raise ValueError(f"Column '{y_col}' not found in perf_df.")

    # Drop any NaN rows for x_col or y_col
    df_plot = perf_df[[x_col, y_col]].dropna()
    if verbose:
        print(f"[info] Plotting 2D distribution from shape={df_plot.shape} after dropping NA in {x_col},{y_col}.")

    # Initialize figure
    fig, ax = plt.subplots(figsize=figsize)

    # Default title
    if title is None:
        title = f"{plot_type.capitalize()} Plot: {y_col} vs. {x_col}"

    # Depending on plot_type, create different plots
    if plot_type.lower() == "density":
        # Two main approaches:
        # (A) scatter + sns.kdeplot
        # (B) scatter + hist2d
        # We show an example with scatter + kde shading.

        # Light scatterplot
        sns.scatterplot(data=df_plot, x=x_col, y=y_col, 
                        color="blue", alpha=0.4, s=15, ax=ax)
        
        # 2D density (kde)
        sns.kdeplot(
            data=df_plot,
            x=x_col,
            y=y_col,
            fill=True,
            cmap="mako",
            alpha=0.6,
            levels=10,
            thresh=0,    # ensures we fill entire range
            ax=ax
        )
        ax.set_title(title)

    elif plot_type.lower() == "heatmap":
        # We can do a 2D histogram -> pivot -> heatmap
        xvals = df_plot[x_col]
        yvals = df_plot[y_col]

        # Using numpy.histogram2d
        h, xedges, yedges = np.histogram2d(xvals, yvals, bins=bins)
        
        # Edges define bin boundaries in x, y
        # Construct coordinates for heatmap
        # Easiest way is to show via sns.heatmap if we pivot, but we'll do an imshow approach

        # Option 1: direct imshow (like pcolormesh)
        # Easiest might be to flip h with .T so y is "top to bottom" in the right orientation
        # But let's do a simpler approach with pcolormesh or ax.imshow
        mesh = ax.imshow(
            h.T, 
            origin="lower",
            aspect="auto",
            extent=[xedges[0], xedges[-1], yedges[0], yedges[-1]],
            cmap="magma"
        )
        fig.colorbar(mesh, ax=ax, label="Count (# genes)")
        ax.set_title(title)
        ax.set_xlabel(x_col)
        ax.set_ylabel(y_col)

    else:
        # fallback => basic scatter
        ax.scatter(df_plot[x_col], df_plot[y_col], alpha=0.4, s=15)
        ax.set_title(f"Scatter only: {y_col} vs. {x_col}")

    ax.set_xlabel(x_col)
    ax.set_ylabel(y_col)

    # Save or show
    if output_path is not None:
        plt.savefig(output_path, dpi=300, bbox_inches="tight")
        if verbose:
            print(f"[info] Saved 2D performance plot to: {output_path}")

    if show:
        plt.show()
    else:
        plt.close(fig)

    return fig, ax


def plot_model_roc(
    position_df: pd.DataFrame,
    model_name: str = "Model",
    score_col: str = "score",
    fold_id_col: Optional[str] = None,
    title: str = "ROC Curve",
    output_path: Optional[str] = None,
    show_plot: bool = True,
    verbose: int = 1
):
    """
    Plot an ROC curve from a 'position dataframe' containing model predictions
    on actual vs. non-actual splice sites.

    Interprets:
      - TP or FN => real splice site => label=1
      - FP or TN => not a real splice site => label=0

    Parameters
    ----------
    position_df : pd.DataFrame
        Must have columns:
          - pred_type ∈ {TP, FP, TN, FN}
          - [optionally] fold_id_col if you want per-fold ROC
          - score_col for predicted probability
    model_name : str
        Name of the model (e.g., SpliceAI, DNABERT).
    score_col : str
        Column name with the model's predicted probability (default "score").
    fold_id_col : str or None
        If present, each fold is plotted separately, plus an average curve.
    title : str
        Plot title.
    output_path : str or None
        If not None, path to save the figure (PDF, PNG, etc.).
    show_plot : bool
        Whether to show via plt.show() if output_path is None.
    verbose : int
        Print progress if > 0.

    Returns
    -------
    None
    """
    # 1) Create a copy and map pred_type -> binary label
    df = position_df.copy()
    # keep only rows with pred_type in [TP, FP, TN, FN]
    df = df[df["pred_type"].isin(["TP", "FP", "TN", "FN"])]

    # If empty after filtering, abort
    if df.empty:
        if verbose:
            print("[warning] No valid pred_type in {TP, FP, TN, FN}. Cannot plot ROC.")
        return

    def map_label(pt):
        # Real site => 1 => (TP or FN)
        # Not real => 0 => (FP or TN)
        return 1 if pt in ("TP", "FN") else 0
    df["y_true"] = df["pred_type"].apply(map_label)
    y_col = "y_true"

    # 2) If we have a fold_id_col, do a per-fold approach
    if fold_id_col is not None and fold_id_col in df.columns:
        folds = df[fold_id_col].unique()
        if verbose:
            print(f"[info] Found fold_id_col='{fold_id_col}' with {len(folds)} folds.")

        tprs = []
        aucs = []
        mean_fpr = np.linspace(0, 1, 100)

        fig = plt.figure(figsize=(8, 6))
        for fold in folds:
            sub = df[df[fold_id_col] == fold]
            if sub.empty:
                continue
            y_true = sub[y_col].values
            y_score = sub[score_col].values

            fpr, tpr, _ = roc_curve(y_true, y_score)
            roc_auc = auc(fpr, tpr)
            aucs.append(roc_auc)

            plt.plot(fpr, tpr, alpha=0.3, label=f"Fold {fold} (AUC={roc_auc:.2f})")

            # Interpolation for average
            interp_tpr = np.interp(mean_fpr, fpr, tpr)
            interp_tpr[0] = 0.0
            tprs.append(interp_tpr)

        # plot average
        if len(tprs) > 0:
            mean_tpr = np.mean(tprs, axis=0)
            mean_tpr[-1] = 1.0
            mean_auc = np.mean(aucs)
            std_auc = np.std(aucs)

            plt.plot(mean_fpr, mean_tpr, 'b', lw=2,
                     label=f"Mean ROC (AUC={mean_auc:.2f} ± {std_auc:.2f})")

        plt.plot([0, 1], [0, 1], 'r--', label="Random")
        plt.title(f"{title} - {model_name} [CV folds]")
        plt.xlabel("False Positive Rate")
        plt.ylabel("True Positive Rate")
        plt.legend(loc="best")
        plt.grid()

        if output_path:
            if verbose:
                print(f"[i/o] Saving ROC to: {output_path}")
            plt.savefig(output_path, bbox_inches='tight', dpi=150)
            plt.close(fig)
        else:
            if show_plot:
                plt.show()
            else:
                plt.close(fig)

    else:
        # single global curve
        y_true = df[y_col].values
        y_score = df[score_col].values
        fpr, tpr, _ = roc_curve(y_true, y_score)
        roc_auc = auc(fpr, tpr)

        fig = plt.figure(figsize=(8, 6))
        plt.plot(fpr, tpr, label=f"ROC (AUC={roc_auc:.2f})", lw=2)
        plt.plot([0, 1], [0, 1], 'r--', label="Random")
        plt.title(f"{title} - {model_name}")
        plt.xlabel("False Positive Rate")
        plt.ylabel("True Positive Rate")
        plt.legend(loc="best")
        plt.grid()

        if output_path:
            if verbose:
                print(f"[i/o] Saving ROC to: {output_path}")
            plt.savefig(output_path, bbox_inches='tight', dpi=150)
            plt.close(fig)
        else:
            if show_plot:
                plt.show()
            else:
                plt.close(fig)


def plot_spliceai_roc(
    position_df: pd.DataFrame,
    score_col: str = "score",
    fold_id_col: Optional[str] = None,
    title: str = "SpliceAI ROC: 'Is it a real splice site?'",
    output_path: Optional[str] = None,
    show_plot: bool = True,
    verbose: int = 1
):
    """
    Plot an ROC curve from a 'position dataframe' containing SpliceAI predictions
    on actual vs. non-actual splice sites.

    Interprets:
      - TP or FN => real splice site => label=1
      - FP or TN => not a real splice site => label=0

    'score_col' is SpliceAI's probability that it's a real site.
    We then compute and plot ROC.

    Parameters
    ----------
    position_df : pd.DataFrame
        Must have columns:
          - pred_type ∈ {TP, FP, TN, FN}
          - [optionally] fold_id_col if you want per-fold ROC
          - score_col for predicted probability
    score_col : str
        Column name with SpliceAI's predicted probability (default "score").
    fold_id_col : str or None
        If present, each fold is plotted separately, plus an average curve.
    title : str
        Plot title.
    output_path : str or None
        If not None, path to save the figure (PDF, PNG, etc.).
    show_plot : bool
        Whether to show via plt.show() if output_path is None.
    verbose : int
        Print progress if > 0.

    Returns
    -------
    None
    """
    from sklearn.metrics import roc_curve, auc

    if isinstance(position_df, pl.DataFrame):
        position_df = position_df.to_pandas()

    # 1) Create a copy and map pred_type -> binary label
    df = position_df.copy()
    # keep only rows with pred_type in [TP, FP, TN, FN]
    df = df[df["pred_type"].isin(["TP","FP","TN","FN"])]

    # If empty after filtering, abort
    if df.empty:
        if verbose:
            print("[warning] No valid pred_type in {TP, FP, TN, FN}. Cannot plot ROC.")
        return

    def map_label(pt):
        # Real site => 1 => (TP or FN)
        # Not real => 0 => (FP or TN)
        return 1 if pt in ("TP","FN") else 0
    df["y_true"] = df["pred_type"].apply(map_label)
    y_col = "y_true"

    # 2) If we have a fold_id_col, do a per-fold approach
    if fold_id_col is not None and fold_id_col in df.columns:
        folds = df[fold_id_col].unique()
        if verbose:
            print(f"[info] Found fold_id_col='{fold_id_col}' with {len(folds)} folds.")

        tprs = []
        aucs = []
        mean_fpr = np.linspace(0,1,100)

        fig = plt.figure(figsize=(8,6))
        for fold in folds:
            sub = df[df[fold_id_col]==fold]
            if sub.empty:
                continue
            y_true = sub[y_col].values
            y_score = sub[score_col].values

            fpr, tpr, _ = roc_curve(y_true, y_score)
            roc_auc = auc(fpr, tpr)
            aucs.append(roc_auc)

            plt.plot(fpr, tpr, alpha=0.3, label=f"Fold {fold} (AUC={roc_auc:.2f})")

            # Interpolation for average
            interp_tpr = np.interp(mean_fpr, fpr, tpr)
            interp_tpr[0] = 0.0
            tprs.append(interp_tpr)

        # plot average
        if len(tprs) > 0:
            mean_tpr = np.mean(tprs, axis=0)
            mean_tpr[-1] = 1.0
            mean_auc = np.mean(aucs)
            std_auc = np.std(aucs)

            plt.plot(mean_fpr, mean_tpr, 'b', lw=2,
                     label=f"Mean ROC (AUC={mean_auc:.2f} ± {std_auc:.2f})")

        plt.plot([0,1],[0,1],'r--', label="Random")
        plt.title(title + " [CV folds]")
        plt.xlabel("False Positive Rate")
        plt.ylabel("True Positive Rate")
        plt.legend(loc="best")
        plt.grid()

        if output_path:
            if verbose:
                print(f"[i/o] Saving ROC to: {output_path}")
            plt.savefig(output_path, bbox_inches='tight', dpi=150)
            plt.close(fig)
        else:
            if show_plot:
                plt.show()
            else:
                plt.close(fig)

    else:
        # single global curve
        y_true = df[y_col].values
        y_score = df[score_col].values
        fpr, tpr, _ = roc_curve(y_true, y_score)
        roc_auc = auc(fpr, tpr)

        fig = plt.figure(figsize=(8,6))
        plt.plot(fpr, tpr, label=f"ROC (AUC={roc_auc:.2f})", lw=2)
        plt.plot([0,1],[0,1],'r--', label="Random")
        plt.title(title)
        plt.xlabel("False Positive Rate")
        plt.ylabel("True Positive Rate")
        plt.legend(loc="best")
        plt.grid()

        if output_path:
            if verbose:
                print(f"[i/o] Saving ROC to: {output_path}")
            plt.savefig(output_path, bbox_inches='tight', dpi=150)
            plt.close(fig)
        else:
            if show_plot:
                plt.show()
            else:
                plt.close(fig)


def plot_model_pr(
    position_df: pd.DataFrame,
    model_name: str = "Model",
    score_col: str = "score",
    fold_id_col: Optional[str] = None,
    title: str = "Precision-Recall Curve",
    output_path: Optional[str] = None,
    show_plot: bool = True,
    verbose: int = 1
):
    """
    Plot a Precision-Recall curve from a 'position dataframe' with model predictions.
    
    Interprets:
      - TP or FN => real site => label=1
      - FP or TN => not a real site => label=0

    Parameters
    ----------
    position_df : pd.DataFrame
        Must have columns:
          - pred_type ∈ {TP, FP, TN, FN}
          - [optionally] fold_id_col
          - score_col
    model_name : str
        Name of the model (e.g., SpliceAI, DNABERT).
    score_col : str
        Column with model probabilities (default "score").
    fold_id_col : str or None
        If present, plot each fold's PR plus average.
    title : str
        Plot title.
    output_path : str or None
        If provided, path to save figure.
    show_plot : bool
        If True, plt.show() if no output_path.
    verbose : int
        Print progress if > 0.

    Returns
    -------
    None
    """
    if isinstance(position_df, pl.DataFrame):
        position_df = position_df.to_pandas()

    df = position_df.copy()
    df = df[df["pred_type"].isin(["TP", "FP", "TN", "FN"])]
    if df.empty:
        if verbose:
            print("[warning] No valid pred_type in {TP, FP, TN, FN}. Cannot plot PR.")
        return

    def map_label(pt):
        # real site => 1 => TP or FN
        # not real => 0 => FP or TN
        return 1 if pt in ("TP", "FN") else 0
    df["y_true"] = df["pred_type"].apply(map_label)

    if fold_id_col is not None and fold_id_col in df.columns:
        folds = df[fold_id_col].unique()
        if verbose:
            print(f"[info] Found fold_id_col='{fold_id_col}' with {len(folds)} folds.")

        plt.figure(figsize=(8, 6))
        mean_recall = np.linspace(0, 1, 100)
        pr_curves = []
        aucs = []

        for fold in folds:
            sub = df[df[fold_id_col] == fold]
            if sub.empty:
                continue
            y_true = sub["y_true"].values
            y_score = sub[score_col].values

            prec, rec, _ = precision_recall_curve(y_true, y_score)
            # area under PR
            pr_auc = auc(rec, prec)
            aucs.append(pr_auc)

            # ensure recall is ascending
            if rec[0] > rec[-1]:
                prec, rec = prec[::-1], rec[::-1]

            plt.plot(rec, prec, alpha=0.3, label=f"Fold {fold} (AUC={pr_auc:.2f})")

            # Interpolate at mean_recall
            prec_interp = np.interp(mean_recall, rec, prec)
            pr_curves.append(prec_interp)

        # mean curve
        if len(pr_curves) > 0:
            mean_prec = np.mean(pr_curves, axis=0)
            mean_auc = np.mean(aucs)
            std_auc = np.std(aucs)

            plt.plot(mean_recall, mean_prec, color='b',
                     label=f"Mean PR (AUC={mean_auc:.2f} ± {std_auc:.2f})",
                     lw=2)

        plt.xlim([0, 1])
        plt.ylim([0, 1])
        plt.xlabel("Recall")
        plt.ylabel("Precision")
        plt.title(f"{title} - {model_name} [CV folds]")
        plt.legend(loc="best")
        plt.grid()

        if output_path:
            if verbose:
                print(f"[i/o] Saving PR to: {output_path}")
            plt.savefig(output_path, bbox_inches='tight', dpi=150)
            plt.close()
            print(f"[i/o] Saved PR to: {output_path}")
        else:
            if show_plot:
                plt.show()
            else:
                plt.close()

    else:
        # single global PR
        y_true = df["y_true"].values
        y_score = df[score_col].values

        prec, rec, _ = precision_recall_curve(y_true, y_score)
        pr_auc = auc(rec, prec)

        plt.figure(figsize=(8, 6))
        plt.plot(rec, prec, label=f"PR (AUC={pr_auc:.2f})", lw=2)
        plt.xlim([0, 1])
        plt.ylim([0, 1])
        plt.xlabel("Recall")
        plt.ylabel("Precision")
        plt.title(f"{title} - {model_name}")
        plt.legend(loc="best")
        plt.grid()

        if output_path:
            if verbose:
                print(f"[i/o] Saving PR to: {output_path}")
            plt.savefig(output_path, bbox_inches='tight', dpi=150)
            plt.close()
            print(f"[i/o] Saved PR to: {output_path}")
        else:
            if show_plot:
                plt.show()
            else:
                plt.close()


def plot_spliceai_pr(
    position_df: pd.DataFrame,
    score_col: str = "score",
    fold_id_col: Optional[str] = None,
    title: str = "SpliceAI PR: 'Is it a real splice site?'",
    output_path: Optional[str] = None,
    show_plot: bool = True,
    verbose: int = 1
):
    """
    Plot a Precision-Recall curve from a 'position dataframe' with SpliceAI predictions.
    
    Interprets:
      - TP or FN => real site => label=1
      - FP or TN => not a real site => label=0

    'score_col' is the probability that it's a real site.
    If 'fold_id_col' is present, we do a fold-by-fold approach.

    Parameters
    ----------
    position_df : pd.DataFrame
        Must have columns:
          - pred_type ∈ {TP, FP, TN, FN}
          - [optionally] fold_id_col
          - score_col
    score_col : str
        Column with SpliceAI probabilities (default "score").
    fold_id_col : str or None
        If present, plot each fold's PR plus average.
    title : str
        Plot title.
    output_path : str or None
        If provided, path to save figure.
    show_plot : bool
        If True, plt.show() if no output_path.
    verbose : int
        Print progress if > 0.

    Returns
    -------
    None
    """
    from sklearn.metrics import precision_recall_curve

    if isinstance(position_df, pl.DataFrame):
        position_df = position_df.to_pandas()

    df = position_df.copy()
    df = df[df["pred_type"].isin(["TP","FP","TN","FN"])]
    if df.empty:
        if verbose:
            print("[warning] No valid pred_type in {TP,FP,TN,FN}. Cannot plot PR.")
        return

    def map_label(pt):
        # real site => 1 => TP or FN
        # not real => 0 => FP or TN
        return 1 if pt in ("TP","FN") else 0
    df["y_true"] = df["pred_type"].apply(map_label)

    if fold_id_col is not None and fold_id_col in df.columns:
        folds = df[fold_id_col].unique()
        if verbose:
            print(f"[info] Found fold_id_col='{fold_id_col}' with {len(folds)} folds.")

        plt.figure(figsize=(8,6))
        mean_recall = np.linspace(0,1,100)
        pr_curves = []
        aucs = []

        for fold in folds:
            sub = df[df[fold_id_col]==fold]
            if sub.empty:
                continue
            y_true = sub["y_true"].values
            y_score = sub[score_col].values

            prec, rec, _ = precision_recall_curve(y_true, y_score)
            # area under PR
            pr_auc = auc(rec, prec)
            aucs.append(pr_auc)

            # ensure recall is ascending
            if rec[0] > rec[-1]:
                prec, rec = prec[::-1], rec[::-1]

            plt.plot(rec, prec, alpha=0.3, label=f"Fold {fold} (AUC={pr_auc:.2f})")

            # Interpolate at mean_recall
            prec_interp = np.interp(mean_recall, rec, prec)
            pr_curves.append(prec_interp)

        # mean curve
        if len(pr_curves) > 0:
            mean_prec = np.mean(pr_curves, axis=0)
            mean_auc = np.mean(aucs)
            std_auc = np.std(aucs)

            plt.plot(mean_recall, mean_prec, color='b',
                     label=f"Mean PR (AUC={mean_auc:.2f} ± {std_auc:.2f})",
                     lw=2)

        plt.xlim([0,1])
        plt.ylim([0,1])
        plt.xlabel("Recall")
        plt.ylabel("Precision")
        plt.title(title + " [CV folds]")
        plt.legend(loc="best")
        plt.grid()

        if output_path:
            if verbose:
                print(f"[i/o] Saving PR to: {output_path}")
            plt.savefig(output_path, bbox_inches='tight', dpi=150)
            plt.close()
            print(f"[i/o] Saved PR to: {output_path}")
        else:
            if show_plot:
                plt.show()
            else:
                plt.close()

    else:
        # single global PR
        y_true = df["y_true"].values
        y_score = df[score_col].values

        prec, rec, _ = precision_recall_curve(y_true, y_score)
        pr_auc = auc(rec, prec)

        plt.figure(figsize=(8,6))
        plt.plot(rec, prec, label=f"PR (AUC={pr_auc:.2f})", lw=2)
        plt.xlim([0,1])
        plt.ylim([0,1])
        plt.xlabel("Recall")
        plt.ylabel("Precision")
        plt.title(title)
        plt.legend(loc="best")
        plt.grid()

        if output_path:
            if verbose:
                print(f"[i/o] Saving PR to: {output_path}")
            plt.savefig(output_path, bbox_inches='tight', dpi=150)
            plt.close()
            print(f"[i/o] Saved PR to: {output_path}")
        else:
            if show_plot:
                plt.show()
            else:
                plt.close()


############################################


def demo_analyze_hard_genes(**kargs): 

    eval_dir = kargs.get('eval_dir', PerformanceAnalyzer.eval_dir)
    separator = kargs.get('separator', kargs.get('sep', '\t'))
    error_type = kargs.get('error_type', 'FP')

    mefd = ModelEvaluationFileHandler(eval_dir, separator=separator)  # Initialize ModelEvaluationFileHandler

    print_emphasized("[i/o] Loading performance data...")
    full_performance_df = mefd.load_performance_df(aggregated=True)  # Load the performance dataframe

    pa = PerformanceAnalyzer()

    print_emphasized(f"[workflow] Analyzing 'hard' genes (for which SpliceAI didn't do well) ...")
    pa.analyze_hard_genes(
        full_performance_df, 
        metric='f1_score', 
        error_type=error_type, 
        n_genes=500, save=True)

    
def demo_plot_gene_performance_2d(**kargs): 

    eval_dir = kargs.get('eval_dir', PerformanceAnalyzer.eval_dir)
    separator = kargs.get('separator', kargs.get('sep', '\t'))

    mefd = ModelEvaluationFileHandler(eval_dir, separator=separator)  # Initialize ModelEvaluationFileHandler

    print_emphasized("[i/o] Loading performance data...")
    full_performance_df = mefd.load_performance_df(aggregated=True)  # Load the performance dataframe

    analyze_dataframe_properties(full_performance_df)

    pa = PerformanceAnalyzer()

    # Plot the density of the performance data
    output_path = os.path.join(pa.output_dir, "gene_performance_2d_density.pdf")
    plot_gene_performance_2d(
        full_performance_df, 
        x_col="recall", 
        y_col="precision", 
        plot_type="density", 
        bins=50, 
        figsize=(8,6), 
        output_path=output_path
    )

    # Plot the heatmap of the performance data  
    output_path = os.path.join(pa.output_dir, "gene_performance_2d_heatmap.pdf")
    plot_gene_performance_2d(
        full_performance_df, 
        x_col="recall", 
        y_col="precision", 
        plot_type="heatmap", 
        bins=50, 
        output_path=output_path
    )


def demo_plot_position_df_roc(**kargs): 
    from .extract_genomic_features import SpliceAnalyzer
    from meta_spliceai.splice_engine.run_spliceai_workflow import retrieve_predicted_splice_sites

    separator = kargs.get('separator', kargs.get('sep', '\t'))
    splice_site_predictor = 'spliceai'
    local_dir = kargs.get('local_dir', '/path/to/meta-spliceai/data/ensembl/')
    src_dir = kargs.get('src_dir', '/path/to/meta-spliceai/data/ensembl/')
    # eval_dir = \
    #     kargs.get(
    #         'eval_dir', 
    #         os.path.join(src_dir, f"{splice_site_predictor}_eval"))

    # /path/to/meta-spliceai/data/ensembl/spliceai_eval/full_splice_positions_test.tsv
    # mefd = ModelEvaluationFileHandler(eval_dir, separator=separator)  # Initialize ModelEvaluationFileHandler

    print_emphasized(f"Retrieving predicted splice sites ...")
    position_df = retrieve_predicted_splice_sites(local_dir)
    print_with_indent(f"Column: {list(position_df.columns)}", indent_level=1)

    pa = PerformanceAnalyzer()

    # Find unique values in the 'pred_type' column
    unique_pred_types = position_df['pred_type'].unique()
    print_with_indent(f"Unique values in 'pred_type' column: {unique_pred_types}", indent_level=1)

    # sa = SpliceAnalyzer()
    # print_with_indent(f"Splice sites file: {sa.path_to_splice_sites}", indent_level=1)

    # Plot the ROC curve for the position dataframe
    output_path = os.path.join(pa.output_dir, "position_df_roc.pdf")
    plot_spliceai_roc(
        position_df=position_df, 
        score_col="score",
        fold_id_col=None, 
        title="SpliceAI ROC: 'Is it a real splice site?'",
        output_path=output_path,
        show_plot=False, 
        verbose=1
    )

    # Plot the PR curve for the position dataframe
    output_path = os.path.join(pa.output_dir, "position_df_pr.pdf")
    plot_spliceai_pr(
        position_df=position_df, 
        score_col="score",
        fold_id_col=None, 
        title="SpliceAI PR: 'Is it a real splice site?'",
        output_path=output_path,
        show_plot=False, 
        verbose=1
    )




def test(): 

    # print_emphasized("Demo: Analyze Performance Profile")
    # analyze_performance_profile(
    #     # performance_df=df, 
    #     metric_to_plot='f1_score', 
    #     save_plot=True, 
    #     save_stats=True, 
    #     plot_file='precision_distribution.pdf', 
    #     verbose=True
    # )

    # for metric in ['precision', 'recall', 'specificity', 'f1_score', 'fpr', 'fnr']:
    #     analyze_performance_profile(
    #         # performance_df=df, 
    #         metric_to_plot=metric, 
    #         save_plot=True, 
    #         save_stats=False, 
    #         plot_file=f"{metric}_distribution.pdf", 
    #         verbose=True
    #     )

    # print_emphasized("Demo: Sort Performance by Metric and Error Type")
    # demo_sort_performance_by_metric_and_error_type()

    print_emphasized("Demo: Analyzing Hard Genes (for which SpliceAI didn't do well) ...")
    # analyze_hard_genes(metric='f1_score', error_type='FP', n_genes=100)

    # demo_analyze_hard_genes()

    print_emphasized("Demo: Plotting Gene Performance 2D ...")
    # demo_plot_gene_performance_2d()

    print_emphasized("Demo: Plotting Position DataFrame ROC ...")
    demo_plot_position_df_roc()
    

if __name__ == "__main__":
   test()    
