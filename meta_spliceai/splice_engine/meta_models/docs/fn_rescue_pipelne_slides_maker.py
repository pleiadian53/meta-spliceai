from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)

add_title_slide("Meta‑Model Feature Suite", "MetaSpliceAI")

add_content_slide("Raw SpliceAI Outputs", [
"Donor (p_d), Acceptor (p_a), Neither (p_n)",
"Probability triplet used as primary features",
"Sum to 1 per nucleotide"
])

add_content_slide("Phase 1 – Basic Derived Features", [
"Relative donor = p_d / (p_d+p_a)",
"Splice prob = p_d + p_a",
"Δ donor–acceptor = (p_d−p_a)/(p_d+p_a+ε)",
"Log‑odds donor = log((p_d+ε)/(p_a+ε))",
"Entropy = −Σ p_i log₂ p_i"
])

add_content_slide("Phase 2 – Context Window (±2 nt)", [
"Scores m2 m1 0 p1 p2 extracted per splice type",
"Zero‑padded at transcript ends"
])

add_content_slide("Context‑Differential Metrics", [
"diff_m1 = score₀ − score_{-1}",
"surge_ratio = score₀ / (score_{-1}+score_{+1})",
"neighbor_mean, signal_strength",
"local_peak flag"
])

add_content_slide("Signal‑Processing Metrics", [
"Second derivative ≈ score_{-1} −2·score₀ + score_{+1}",
"Asymmetry = mean(upstream) − mean(downstream)"
])

add_content_slide("Phase 3 – Cross‑Type Features", [
"peak_ratio = donor_peak / acceptor_peak",
"type_signal_difference = donor_strength − acceptor_strength"
])

add_content_slide("Why They Help", [
"Shape metrics capture true splice geometry",
"Entropy highlights ambiguous calls (likely FP/FN)",
"Context ratios rescue weak but canonical FNs",
"Cross‑type ratios suppress dual‑channel decoys"
])

add_content_slide("Empirical Gains", [
"+4‑6 % recall uplift on held‑out test genes",
"Precision unchanged after probability floor 0.005"
])

file_path = "/mnt/data/meta_model_feature_deck.pptx"
prs.save(file_path)
file_path
